#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
OUTPUT_PATH = ROOT / "thesis_presentation.pptx"

LOGO_PATH = ROOT / "Thesis Latex" / "images" / "logos" / "uva_logo.jpg"
PIPELINE_PATH = ROOT / "Thesis Latex" / "images" / "pipeline_architecture.png"
FIGURES = {
    "base_model_comparison": ROOT / "outputs" / "figures" / "base_model_comparison.png",
    "error_breakdown": ROOT / "outputs" / "figures" / "error_breakdown.png",
    "baseline_first_run": ROOT / "outputs" / "figures" / "baseline_first_run.png",
    "decomp_ablation": ROOT / "outputs" / "figures" / "decomp_ablation.png",
    "lenient_match_27b": ROOT / "outputs" / "figures" / "lenient_match_27b.png",
    "sft_results": ROOT / "outputs" / "figures" / "sft_results.png",
    "cache_sweep": ROOT / "outputs" / "figures" / "cache_sweep.png",
}

RED = RGBColor(0xBC, 0x00, 0x31)
DARK = RGBColor(0x22, 0x22, 0x22)
DARK_SECTION = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF8, 0xF8)
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)
ROW_GREY = RGBColor(0xF0, 0xF0, 0xF0)
BLUE_FILL = RGBColor(0xEA, 0xF4, 0xFB)
BLUE_BORDER = RGBColor(0x6E, 0x8A, 0x9F)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def all_assets_exist() -> None:
    required = [LOGO_PATH, PIPELINE_PATH, *FIGURES.values()]
    missing = [str(path) for path in required if not path.exists()]
    if missing:
        raise FileNotFoundError("Missing required assets:\n" + "\n".join(missing))


def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_shape_fill(shape, color: RGBColor) -> None:
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_shape_line(shape, color: RGBColor, width_pt: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    font_size: float = 18,
    color: RGBColor = DARK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = "Calibri",
    vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margins: tuple[float, float, float, float] = (0.06, 0.04, 0.06, 0.04),
) :
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = vertical_anchor
    tf.margin_left = Inches(margins[0])
    tf.margin_top = Inches(margins[1])
    tf.margin_right = Inches(margins[2])
    tf.margin_bottom = Inches(margins[3])
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    return box, tf


def add_bullets(
    slide,
    items: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: float = 17,
    text_color: RGBColor = DARK,
    bullet_color: RGBColor = RED,
    bullet_char: str = "•",
    space_after: float = 8,
) :
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.clear()

    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        p.line_spacing = 1.12

        bullet_run = p.add_run()
        bullet_run.text = f"{bullet_char} "
        bullet_font = bullet_run.font
        bullet_font.name = "Calibri"
        bullet_font.size = Pt(font_size)
        bullet_font.bold = True
        bullet_font.color.rgb = bullet_color

        text_run = p.add_run()
        text_run.text = item
        text_font = text_run.font
        text_font.name = "Calibri"
        text_font.size = Pt(font_size)
        text_font.color.rgb = text_color
    return box


def add_picture_contain(slide, image_path: Path, x: float, y: float, max_w: float, max_h: float):
    with Image.open(image_path) as img:
        img_w, img_h = img.size
    scale = min(max_w / img_w, max_h / img_h)
    width = img_w * scale
    height = img_h * scale
    return slide.shapes.add_picture(
        str(image_path),
        Inches(x + (max_w - width) / 2),
        Inches(y + (max_h - height) / 2),
        width=Inches(width),
        height=Inches(height),
    )


def add_logo(slide, *, width: float = 0.9, x: float | None = None, y: float | None = None) -> None:
    if x is None:
        x = 13.33 - width - 0.25
    if y is None:
        y = 7.5 - width - 0.18
    slide.shapes.add_picture(str(LOGO_PATH), Inches(x), Inches(y), width=Inches(width))


def add_top_bar(slide, title: str) -> None:
    set_background(slide, LIGHT_BG)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.0))
    set_shape_fill(bar, RED)
    bar.line.fill.background()
    add_textbox(
        slide,
        0.62,
        0.16,
        11.5,
        0.6,
        title,
        font_size=26,
        color=WHITE,
        bold=True,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
        margins=(0, 0, 0, 0),
    )
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.36),
        Inches(1.18),
        Inches(0.08),
        Inches(5.95),
    )
    set_shape_fill(accent, RED)
    accent.line.fill.background()


def add_section_divider(slide, title: str, subtitle: str) -> None:
    set_background(slide, DARK_SECTION)
    add_textbox(
        slide,
        0.8,
        0.45,
        11.7,
        0.55,
        title,
        font_size=28,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        1.0,
        1.15,
        11.3,
        0.45,
        subtitle,
        font_size=17,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )


def _replace_timing(slide, timing_el) -> None:
    for existing in slide._element.findall(f"{{{PML_NS}}}timing"):
        slide._element.remove(existing)
    ext_lst = slide._element.find(f"{{{PML_NS}}}extLst")
    if ext_lst is None:
        slide._element.append(timing_el)
        return
    slide._element.insert(list(slide._element).index(ext_lst), timing_el)


def add_paragraph_animations(slide, shape) -> None:
    para_indices = [idx for idx, para in enumerate(shape.text_frame.paragraphs) if para.text.strip()]
    if len(para_indices) <= 1:
        return

    shape_id = shape.shape_id
    xml_str = """<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
        <p:childTnLst>"""

    cid = 2
    for group_id, para_idx in enumerate(para_indices):
        xml_str += f"""
          <p:par>
            <p:cTn id="{cid}" fill="hold">
              <p:stCondLst>
                <p:cond evt="onClick" delay="0">
                  <p:tn>
                    <p:prevCondLst>
                      <p:cond evt="begin" delay="0"/>
                    </p:prevCondLst>
                  </p:tn>
                </p:cond>
              </p:stCondLst>
              <p:childTnLst>
                <p:par>
                  <p:cTn id="{cid + 1}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" grpId="{group_id}" nodeType="clickEffect">
                    <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                    <p:childTnLst>
                      <p:set>
                        <p:cBhvr>
                          <p:cTn id="{cid + 2}" dur="1" fill="hold"/>
                          <p:tgtEl>
                            <p:spTgt spid="{shape_id}">
                              <p:txEl><p:pRg st="{para_idx}" end="{para_idx}"/></p:txEl>
                            </p:spTgt>
                          </p:tgtEl>
                          <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                        </p:cBhvr>
                        <p:to><p:strVal val="visible"/></p:to>
                      </p:set>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
          </p:par>"""
        cid += 3

    xml_str += """
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
  <p:bldLst>"""
    for group_id in range(len(para_indices)):
        xml_str += f'<p:bldP spid="{shape_id}" grpId="{group_id}" uiExpand="1" build="p"/>'
    xml_str += """
  </p:bldLst>
</p:timing>"""
    _replace_timing(slide, etree.fromstring(xml_str))


def add_shape_sequence_animations(slide, shape_groups: list[list]) -> None:
    if not shape_groups:
        return

    xml_str = """<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
        <p:childTnLst>"""

    cid = 2
    inner_id = 200
    for group_id, group in enumerate(shape_groups):
        xml_str += f"""
          <p:par>
            <p:cTn id="{cid}" fill="hold">
              <p:stCondLst>
                <p:cond evt="onClick" delay="0">
                  <p:tn>
                    <p:prevCondLst>
                      <p:cond evt="begin" delay="0"/>
                    </p:prevCondLst>
                  </p:tn>
                </p:cond>
              </p:stCondLst>
              <p:childTnLst>"""
        for shape in group:
            xml_str += f"""
                <p:par>
                  <p:cTn id="{inner_id}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" grpId="{group_id}" nodeType="clickEffect">
                    <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                    <p:childTnLst>
                      <p:set>
                        <p:cBhvr>
                          <p:cTn id="{inner_id + 1}" dur="1" fill="hold"/>
                          <p:tgtEl><p:spTgt spid="{shape.shape_id}"/></p:tgtEl>
                          <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                        </p:cBhvr>
                        <p:to><p:strVal val="visible"/></p:to>
                      </p:set>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>"""
            inner_id += 2
        xml_str += """
              </p:childTnLst>
            </p:cTn>
          </p:par>"""
        cid += 1

    xml_str += """
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""
    _replace_timing(slide, etree.fromstring(xml_str))


def add_notes(slide, text: str) -> None:
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    p = notes_tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = "Calibri"
    font.size = Pt(12)


def make_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RED)
    add_textbox(
        slide,
        0.9,
        0.82,
        11.5,
        1.55,
        "NLP-to-BraneScript Translation\nfor Distributed Scientific Workflows",
        font_size=27,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        3.45,
        2.48,
        6.4,
        0.42,
        "Master's Thesis Defence",
        font_size=20,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        3.0,
        2.92,
        7.3,
        0.36,
        "University of Amsterdam · 2025",
        font_size=16,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    slide.shapes.add_picture(str(LOGO_PATH), Inches(5.42), Inches(4.15), width=Inches(2.5))


def make_problem_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "The Problem")
    bullet_box = add_bullets(
        slide,
        [
            "Brane: federated workflow framework — data stays local, computations travel",
            "Workflows written in BraneScript, a custom domain-specific language (DSL)",
            "Barrier: scientists must learn DSL syntax to run analyses",
            "Goal: translate plain-English intent → working BraneScript workflow",
        ],
        0.82,
        1.44,
        5.35,
        4.6,
        font_size=17,
    )
    label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.52), Inches(1.38), Inches(1.3), Inches(0.42))
    set_shape_fill(label, RED)
    label.line.fill.background()
    add_textbox(
        slide,
        6.58,
        1.43,
        1.18,
        0.28,
        "Example",
        font_size=14,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    prompt_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.25), Inches(1.9), Inches(6.2), Inches(1.18))
    set_shape_fill(prompt_box, BLUE_FILL)
    set_shape_line(prompt_box, BLUE_BORDER, 1.5)
    add_textbox(
        slide,
        6.45,
        2.2,
        5.8,
        0.55,
        "\"Compute the mean age of diabetic patients\"",
        font_size=18,
        color=DARK,
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        8.86,
        3.12,
        0.35,
        0.35,
        "↓",
        font_size=24,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.02), Inches(3.52), Inches(6.6), Inches(1.92))
    set_shape_fill(code_box, BLUE_FILL)
    set_shape_line(code_box, BLUE_BORDER, 1.5)
    add_textbox(
        slide,
        6.28,
        3.78,
        6.05,
        1.36,
        'val result := compute_mean(dataset := ..., column := "age",\nfilter := "diabetic == true");\nprintln(result);',
        font_size=15,
        color=DARK,
        font_name="Calibri",
    )
    add_logo(slide)
    add_paragraph_animations(slide, bullet_box)
    add_notes(
        slide,
        "Frame the problem as a usability gap: Brane is powerful for privacy-preserving workflows, but writing correct BraneScript remains a barrier for domain scientists.",
    )


def make_research_questions_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "Research Questions")
    questions = [
        ("RQ1", "How does model scale affect translation accuracy?"),
        ("RQ2", "What is the effect of retrieval, prompt engineering, and decomposition?"),
        ("RQ3", "Does supervised fine-tuning on domain examples improve accuracy?"),
        ("RQ4", "Can a semantic cache provide reproducibility for equivalent intents?"),
    ]
    y = 1.42
    shape_groups: list[list] = []
    for rq, text in questions:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(y), Inches(11.65), Inches(1.08))
        set_shape_fill(box, WHITE)
        set_shape_line(box, RGBColor(0xD9, 0xD9, 0xD9), 1.0)
        box.text_frame.clear()
        box.text_frame.word_wrap = True
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        box.text_frame.margin_left = Inches(0.42)
        box.text_frame.margin_right = Inches(0.18)
        box.text_frame.margin_top = Inches(0.06)
        box.text_frame.margin_bottom = Inches(0.04)
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        rq_run = p.add_run()
        rq_run.text = f"{rq}  "
        rq_run.font.name = "Calibri"
        rq_run.font.size = Pt(18)
        rq_run.font.bold = True
        rq_run.font.color.rgb = RED
        text_run = p.add_run()
        text_run.text = text
        text_run.font.name = "Calibri"
        text_run.font.size = Pt(18)
        text_run.font.color.rgb = DARK
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.92), Inches(y), Inches(0.13), Inches(1.08))
        set_shape_fill(band, RED)
        band.line.fill.background()
        shape_groups.append([box, band])
        y += 1.31
    add_logo(slide)
    add_shape_sequence_animations(slide, shape_groups)


def make_architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(slide, "System Architecture", "RAG · Decomposition · Generation · Cache")
    add_picture_contain(slide, PIPELINE_PATH, 1.65, 1.95, 10.0, 4.85)
    add_logo(slide, width=0.72, y=6.56)


def make_dataset_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "Dataset & Evaluation Setup")
    bullet_box = add_bullets(
        slide,
        [
            "584 unique intent–BraneScript pairs",
            "7 domain packages: healthcare, genomics, epidemics, statistics, text analysis, data masking, datetime",
            "Split: 497 training / 87 test (stratified)",
            "3 paraphrases per intent → 1,752 cache evaluation queries",
        ],
        0.82,
        1.44,
        5.85,
        4.85,
        font_size=16.5,
    )
    add_textbox(slide, 7.28, 1.65, 4.8, 0.3, "Evaluation pipeline", font_size=16, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    box_specs = [
        (7.1, 2.45, 1.72, "Compile\nrate"),
        (9.22, 2.45, 1.72, "Execution\nrate"),
        (11.34, 2.45, 1.72, "Output\nmatch rate"),
    ]
    for x, y, w, text in box_specs:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.08))
        set_shape_fill(shape, RED)
        set_shape_line(shape, RED, 1.0)
        add_textbox(
            slide,
            x + 0.1,
            y + 0.18,
            w - 0.2,
            0.62,
            text,
            font_size=15,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
    for x in [8.83, 10.95]:
        add_textbox(slide, x, 2.69, 0.3, 0.25, "→", font_size=22, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_logo(slide)
    add_paragraph_animations(slide, bullet_box)


def make_figure_bullets_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    image_box: tuple[float, float, float, float],
    bullets: list[str],
    *,
    bullet_box: tuple[float, float, float, float] = (8.0, 1.7, 4.55, 4.35),
    font_size: float = 16.5,
    notes: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, title)
    add_picture_contain(slide, image_path, *image_box)
    bullet_shape = add_bullets(slide, bullets, *bullet_box, font_size=font_size)
    add_logo(slide)
    add_paragraph_animations(slide, bullet_shape)
    if notes:
        add_notes(slide, notes)


def make_error_breakdown_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "RQ1 — Error Type Breakdown")
    add_picture_contain(slide, FIGURES["error_breakdown"], 2.1, 1.5, 8.8, 4.75)
    add_textbox(
        slide,
        1.85,
        6.02,
        9.65,
        0.32,
        "Error categories across model sizes. 27B largely eliminates compile failures.",
        font_size=14,
        color=DARK,
        align=PP_ALIGN.CENTER,
    )
    add_logo(slide)


def make_lenient_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "RQ2 — Output Match Sensitivity")
    add_picture_contain(slide, FIGURES["lenient_match_27b"], 0.78, 1.52, 6.3, 4.95)
    add_textbox(slide, 7.35, 1.58, 4.45, 0.28, "Why the strict metric misses some successes", font_size=15, color=DARK, bold=True)
    bullet_box = add_bullets(
        slide,
        [
            "Strict metric penalises correct computations with over-annotated output",
            "26/159 failures: correct values, wrong format",
            'Example: "52.27" vs "Mean: 52.27"',
            "Lenient: 49.6% → 66.0% (+16.4% of failures)",
            "All primary results use strict metric",
        ],
        7.28,
        1.95,
        5.12,
        4.8,
        font_size=15,
        space_after=6,
    )
    add_logo(slide)
    add_paragraph_animations(slide, bullet_box)


def make_threats_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "Threats to Validity")
    rows = [
        ("INTERNAL", RED, "Round 1→2 multi-variable confound; SFT test/train overlap"),
        ("EXTERNAL", DARK, "Single DSL · single model family · English-only · single author"),
        ("CONSTRUCT", RED, "Execution rate ≠ correctness; output match = string equality"),
        ("CONCLUSION", DARK, "87-example test set; single trial per configuration"),
    ]
    y = 1.55
    for idx, (label, color, text) in enumerate(rows):
        body_fill = WHITE if idx % 2 == 0 else ROW_GREY
        body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2.35), Inches(y), Inches(9.88), Inches(0.86))
        set_shape_fill(body, body_fill)
        set_shape_line(body, RGBColor(0xE0, 0xE0, 0xE0), 1.0)
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.96), Inches(y), Inches(1.55), Inches(0.86))
        set_shape_fill(tag, color)
        tag.line.fill.background()
        add_textbox(
            slide,
            1.03,
            y + 0.23,
            1.38,
            0.26,
            label,
            font_size=14.5,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(slide, 2.58, y + 0.2, 9.4, 0.34, text, font_size=14.6, color=DARK)
        y += 1.08
    add_logo(slide)


def make_discussion_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "Discussion")
    bullet_box = add_bullets(
        slide,
        [
            "Scale + prompt engineering: strongest lever",
            "~50% output match ceiling is structural (multi-step logic, package selection)",
            "Semantic cache: effective reproducibility layer — errors caught by execution verification",
            "SFT negative result is itself informative: bootstrapped data insufficient",
        ],
        0.9,
        1.62,
        11.15,
        4.7,
        font_size=18,
    )
    add_logo(slide)
    add_paragraph_animations(slide, bullet_box)


def make_future_work_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "Future Work")
    bullet_box = add_bullets(
        slide,
        [
            "Controlled ablation of each component independently",
            "Real scientist users as intent authors (diverse phrasing)",
            "Closed-model comparison (GPT-4, Claude)",
            "Parameter-aware cache embeddings",
            "SFT with human-verified, independently authored training data",
        ],
        0.9,
        1.58,
        11.15,
        4.95,
        font_size=17.5,
    )
    add_logo(slide)
    add_paragraph_animations(slide, bullet_box)


def make_conclusion_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RED)
    add_textbox(
        slide,
        1.0,
        0.88,
        11.3,
        0.98,
        "Plain-English → BraneScript translation is feasible",
        font_size=27,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_bullets(
        slide,
        [
            "27B + RAG + decomposition: 49.6% strict / 66% lenient output match",
            "Prompt engineering as important as model scale",
            "Semantic cache: 99.1% hit rate, 98.3% correct rate",
            "SFT in this configuration does not help — a finding in itself",
        ],
        1.35,
        2.45,
        10.7,
        2.6,
        font_size=17,
        text_color=WHITE,
        bullet_color=WHITE,
        space_after=8,
    )
    slide.shapes.add_picture(str(LOGO_PATH), Inches(5.76), Inches(5.48), width=Inches(0.95))
    add_textbox(
        slide,
        5.18,
        6.52,
        3.0,
        0.32,
        "Thank you",
        font_size=18,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def build_presentation() -> Presentation:
    all_assets_exist()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    make_title_slide(prs)
    make_problem_slide(prs)
    make_research_questions_slide(prs)
    make_architecture_slide(prs)
    make_dataset_slide(prs)
    make_figure_bullets_slide(
        prs,
        "RQ1 — Translation Accuracy Across Model Scales",
        FIGURES["base_model_comparison"],
        (0.72, 1.52, 7.0, 4.9),
        [
            "Largest jump: 9B → 27B",
            "27B achieves 49.6% output match",
            "4B model dominated by undefined variable errors",
            "Scale is the strongest single lever",
        ],
        notes="Highlight the nonlinear improvement with scale: the step from 9B to 27B changes the error profile from compile failures to mostly harder semantic errors.",
    )
    make_error_breakdown_slide(prs)
    make_figure_bullets_slide(
        prs,
        "RQ2 — Effect of Prompt Engineering",
        FIGURES["baseline_first_run"],
        (0.72, 1.52, 7.0, 4.9),
        [
            "Same 27B weights — two rounds",
            "Round 1: 13.5% output match",
            "Round 2: 49.6% output match",
            "+36.1 percentage points",
            "⚠ Multiple variables changed simultaneously",
        ],
        font_size=15.5,
        notes="Stress that prompt engineering matters almost as much as scale here, while also flagging the confound: retrieval, prompting, and decomposition evolved together between rounds.",
    )
    make_figure_bullets_slide(
        prs,
        "RQ2 — Decomposition Ablation",
        FIGURES["decomp_ablation"],
        (0.72, 1.52, 7.0, 4.9),
        [
            "Removing decomposition: −13.2 pp compile",
            "−16.8 pp execution",
            "−10.9 pp output match",
            "Breaking intents into sub-intents is necessary",
        ],
    )
    make_lenient_slide(prs)
    make_figure_bullets_slide(
        prs,
        "RQ3 — Effect of Supervised Fine-Tuning",
        FIGURES["sft_results"],
        (0.72, 1.52, 7.0, 4.9),
        [
            "3-epoch SFT: decline across all metrics",
            "5-epoch SFT: near-total collapse (~0%)",
            "Base model + good prompt outperforms SFT",
            "Training data bootstrapped from same pipeline — likely limits signal",
        ],
        font_size=15.7,
        notes="Present this as a useful negative result: when the supervision comes from the same imperfect pipeline, fine-tuning can reinforce errors instead of adding signal.",
    )
    make_figure_bullets_slide(
        prs,
        "RQ4 — Semantic Cache as Reproducibility Layer",
        FIGURES["cache_sweep"],
        (0.72, 1.52, 7.0, 4.9),
        [
            "At threshold 0.92:",
            "Hit rate: 99.1%",
            "Correct rate: 98.3%",
            "False-positive rate: 0.8%",
            "Threshold = confidence dial",
            "Higher → fresh generation on uncertain matches",
        ],
        font_size=15.4,
        notes="Explain the cache as a reproducibility layer rather than just an optimisation: equivalent intents resolve to the same workflow unless similarity falls below the confidence threshold.",
    )
    make_threats_slide(prs)
    make_discussion_slide(prs)
    make_future_work_slide(prs)
    make_conclusion_slide(prs)
    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(OUTPUT_PATH)
    reopened = Presentation(str(OUTPUT_PATH))
    if len(reopened.slides) != 16:
        raise RuntimeError(f"Expected 16 slides, found {len(reopened.slides)}")
    print(f"Saved presentation to {OUTPUT_PATH}")
    print(f"Verified slide count: {len(reopened.slides)}")


if __name__ == "__main__":
    main()
